from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PURPLE      = RGBColor(0x5B, 0x21, 0xB6)
LIGHT_PUR   = RGBColor(0xC4, 0xA0, 0xD4)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK        = RGBColor(0x1F, 0x17, 0x3A)
GRAY        = RGBColor(0x6B, 0x72, 0x80)
LAVENDER    = RGBColor(0xF3, 0xF0, 0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # blank layout


def bg(slide, color=LAVENDER):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, left, top, width, height, color=WHITE, radius=False):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def txt(slide, text, left, top, width, height,
        size=24, bold=False, color=DARK, align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def purple_bar(slide, label, top):
    box(slide, 0, top, 13.33, 0.55, PURPLE)
    txt(slide, label, 0.3, top + 0.06, 12, 0.45, size=20, bold=True, color=WHITE)


def bullet_box(slide, items, left, top, width, height, bg_color=WHITE):
    box(slide, left, top, width, height, bg_color)
    txb = slide.shapes.add_textbox(Inches(left+0.15), Inches(top+0.1),
                                   Inches(width-0.3), Inches(height-0.2))
    tf = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = item
        run.font.size = Pt(16)
        run.font.color.rgb = DARK


# ─────────────────────────────────────────────────────────────────────────────
# Slide 1 — Title
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, PURPLE)
txt(s, "HerDay 💜", 1, 1.5, 11, 1.4, size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "Eliminating the Invisible Mental Load on Women",
    1, 3.0, 11, 0.8, size=26, color=LIGHT_PUR, align=PP_ALIGN.CENTER)
txt(s, "Powered by Claude AI  •  Built for the Hackathon",
    1, 4.0, 11, 0.6, size=18, color=WHITE, align=PP_ALIGN.CENTER, italic=True)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 2 — The Problem
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
purple_bar(s, "The Problem — What is the Mental Load?", 0)
txt(s, "Women carry an invisible cognitive burden of managing work AND home — simultaneously.",
    0.5, 0.7, 12, 0.6, size=18, italic=True, color=PURPLE)

cards = [
    ("🧠 Cognitive Overload",   "Remembering what to cook, what to buy, which task is due — all at once"),
    ("📋 Task Switching",        "Jumping between work priorities and home responsibilities breaks focus"),
    ("🍳 Meal Planning Fatigue", "Deciding 3 meals a day × 365 days = 1,095 decisions per year"),
    ("🛒 Grocery Mental Map",    "Tracking what's in the fridge, what's needed, what was rejected before"),
    ("⏰ Time Estimation",       "Guessing how long tasks take while juggling meetings and energy levels"),
]
for i, (title, desc) in enumerate(cards):
    col = i % 2
    row = i // 2
    bx = box(s, 0.4 + col * 6.5, 1.5 + row * 1.35, 6.1, 1.2, WHITE)
    txt(s, title, 0.55 + col * 6.5, 1.55 + row * 1.35, 5.8, 0.4, size=15, bold=True, color=PURPLE)
    txt(s, desc,  0.55 + col * 6.5, 1.95 + row * 1.35, 5.8, 0.6, size=13, color=GRAY)

txt(s, "🔑  Result: Decision fatigue, burnout, and the sense that the load is never done.",
    0.4, 6.5, 12.5, 0.6, size=15, bold=True, color=PURPLE)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 3 — The Solution
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
purple_bar(s, "HerDay — Your AI Life-Load Agent", 0)
txt(s, "One conversation. Every domain. Fully coordinated.",
    0.5, 0.65, 12, 0.5, size=18, italic=True, color=PURPLE)

features = [
    ("💬", "Natural Language First",    "Just tell HerDay how you feel — it figures out the rest"),
    ("📋", "Unified Task View",          "Work tasks + home tasks + meetings in one AI-prioritised list"),
    ("🍽", "Smart Meal Planning",        "3 meals planned around your energy, fridge, dietary needs & BMI"),
    ("🛒", "Live Grocery List",          "Auto-built from meal plan — knows what's already in your fridge"),
    ("📊", "Weekly Insights",            "Tracks mental load saved in hours every week"),
    ("🔒", "Private & Secure",           "Per-user data isolation — nothing shared, everything local"),
]
for i, (icon, title, desc) in enumerate(features):
    col = i % 3
    row = i // 3
    box(s, 0.3 + col * 4.3, 1.3 + row * 1.5, 4.0, 1.35, WHITE)
    txt(s, icon,  0.45 + col * 4.3, 1.35 + row * 1.5, 0.6, 0.5, size=22)
    txt(s, title, 1.0  + col * 4.3, 1.35 + row * 1.5, 3.2, 0.45, size=14, bold=True, color=PURPLE)
    txt(s, desc,  0.45 + col * 4.3, 1.8  + row * 1.5, 3.7, 0.6,  size=12, color=GRAY)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 4 — Agent, not a Chatbot
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
purple_bar(s, "Why HerDay is an Agent — Not Just a Chatbot", 0)

box(s, 0.3, 0.65, 5.9, 5.9, RGBColor(0xFF, 0xEE, 0xEE))
txt(s, "❌  A Chatbot", 0.5, 0.7, 5.5, 0.5, size=18, bold=True, color=RGBColor(0xB9, 0x1C, 0x1C))
for i, line in enumerate([
    "• Answers questions",
    "• One domain at a time",
    "• No memory of preferences",
    "• No actions taken",
    "• Needs exact commands",
]):
    txt(s, line, 0.5, 1.3 + i * 0.55, 5.5, 0.5, size=15, color=DARK)

box(s, 6.8, 0.65, 6.1, 5.9, RGBColor(0xF0, 0xFF, 0xF4))
txt(s, "✅  HerDay Agent", 7.0, 0.7, 5.7, 0.5, size=18, bold=True, color=RGBColor(0x06, 0x5F, 0x46))
for i, line in enumerate([
    "• Takes actions — saves meals, creates tasks",
    "• Coordinates work + home + meals at once",
    "• Remembers rejected meals, past energy",
    "• Proactively builds plans without being asked",
    "• Understands natural language & intent",
]):
    txt(s, line, 7.0, 1.3 + i * 0.55, 5.7, 0.5, size=15, color=DARK)

txt(s, "\"I'm exhausted\"  →  detects low energy  →  adapts meals + tasks  →  saves to DB",
    0.3, 6.3, 12.7, 0.55, size=15, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 5 — AI Moments
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
purple_bar(s, "The AI Moments — Where Claude Makes the Difference", 0)

moments = [
    ("🧠  Energy from Language",
     "User types 'I'm drained' → Claude detects low energy → entire plan adapts. No dropdown, no button."),
    ("🍳  Context-Aware Meals",
     "Claude considers fridge items + BMI + dietary preference + cook time + energy — all in one reasoning step."),
    ("🚫  Rejection Memory",
     "Rejected meals are stored. Claude never suggests them again. It learns taste preferences over time."),
    ("🛒  Fridge Intelligence",
     "Grocery list = meal ingredients MINUS fridge items. Claude knows what you already have."),
    ("⚖️  Cross-Domain Reasoning",
     "One message plans work tasks, home tasks, 3 meals and a shopping list simultaneously."),
]
for i, (title, desc) in enumerate(moments):
    box(s, 0.3, 0.65 + i * 1.15, 12.7, 1.05, WHITE)
    txt(s, title, 0.5,  0.68 + i * 1.15, 4.5, 0.45, size=14, bold=True, color=PURPLE)
    txt(s, desc,  0.5,  1.05 + i * 1.15, 12.2, 0.5, size=13, color=DARK)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 6 — How It Works
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
purple_bar(s, "How It Works — Architecture", 0)

layers = [
    ("💬  User",          "Talks to HerDay in natural language via chat",           WHITE),
    ("🤖  Claude AI",     "Reasons across context — energy, tasks, meals, fridge",  RGBColor(0xED, 0xE9, 0xFE)),
    ("⚙️  Agent Layer",   "Extracts intent, calls functions, takes actions",         RGBColor(0xDD, 0xD6, 0xFE)),
    ("🗄  SQLite DB",     "Stores tasks, meals, logs — per user, fully private",     RGBColor(0xC4, 0xB5, 0xFD)),
    ("📊  Streamlit UI",  "Chat-first interface, lavender theme, real-time updates", WHITE),
]
for i, (layer, desc, color) in enumerate(layers):
    box(s, 1.5, 0.75 + i * 1.1, 10.3, 0.9, color)
    txt(s, layer, 1.7,  0.78 + i * 1.1, 3.5, 0.45, size=15, bold=True, color=PURPLE)
    txt(s, desc,  5.4,  0.78 + i * 1.1, 6.2, 0.45, size=13, color=DARK)
    if i < len(layers) - 1:
        txt(s, "▼", 6.3, 1.55 + i * 1.1, 0.5, 0.3, size=14, color=PURPLE, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 7 — Impact
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
purple_bar(s, "Impact — Mental Load Reduced", 0)

metrics = [
    ("45–55 mins",  "Saved per day\non planning"),
    ("3 meals",     "Planned in\none conversation"),
    ("0",           "Grocery items\nforgotten"),
    ("1 message",   "To coordinate\nyour entire day"),
]
for i, (number, label) in enumerate(metrics):
    box(s, 0.4 + i * 3.15, 0.7, 2.9, 2.2, PURPLE)
    txt(s, number, 0.4 + i * 3.15, 0.85, 2.9, 1.0, size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, label,  0.4 + i * 3.15, 1.8,  2.9, 0.8, size=13, color=LIGHT_PUR, align=PP_ALIGN.CENTER)

before_after = [
    ("Before HerDay",  ["Open 4 separate apps", "Manually plan 3 meals", "Remember fridge contents",
                        "Decide task priorities alone", "Forget grocery items", "Feel overwhelmed"]),
    ("After HerDay",   ["One conversation", "AI plans all 3 meals", "Fridge awareness built-in",
                        "AI prioritises by energy", "Smart grocery list", "Mental load lifted"]),
]
for col, (title, items) in enumerate(before_after):
    color = RGBColor(0xFF, 0xEE, 0xEE) if col == 0 else RGBColor(0xF0, 0xFF, 0xF4)
    tcolor = RGBColor(0xB9, 0x1C, 0x1C) if col == 0 else RGBColor(0x06, 0x5F, 0x46)
    box(s, 0.4 + col * 6.5, 3.1, 6.1, 3.8, color)
    txt(s, ("❌  " if col == 0 else "✅  ") + title,
        0.6 + col * 6.5, 3.15, 5.7, 0.5, size=16, bold=True, color=tcolor)
    for j, item in enumerate(items):
        txt(s, ("• " if col == 0 else "✓ ") + item,
            0.6 + col * 6.5, 3.7 + j * 0.47, 5.7, 0.45, size=13, color=DARK)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 8 — Roadmap
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
purple_bar(s, "What's Next — Roadmap", 0)

roadmap = [
    ("✅  Done",         ["Login & per-user data privacy", "AI meal planning with BMI awareness",
                          "Unified task management", "Live grocery list", "Energy-aware scheduling",
                          "Weekly insights dashboard"]),
    ("🔜  Next Sprint",  ["Google Calendar integration", "Task capture from natural language",
                          "Voice input via Whisper API", "Push notifications — morning briefing"]),
    ("🚀  Future",       ["Adaptive learning from behaviour", "Multi-user family support",
                          "PostgreSQL for scale", "Mobile-responsive UI"]),
]
for i, (phase, items) in enumerate(roadmap):
    box(s, 0.3 + i * 4.3, 0.7, 4.0, 6.3, WHITE)
    txt(s, phase, 0.5 + i * 4.3, 0.75, 3.6, 0.5, size=15, bold=True, color=PURPLE)
    for j, item in enumerate(items):
        txt(s, "• " + item, 0.5 + i * 4.3, 1.35 + j * 0.72, 3.6, 0.65, size=13, color=DARK)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 9 — Close
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, PURPLE)
txt(s, "HerDay 💜", 1, 1.5, 11, 1.2, size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "Plan less. Live more.", 1, 2.9, 11, 0.7, size=28, color=LIGHT_PUR, align=PP_ALIGN.CENTER, italic=True)
txt(s, "Because the mental load shouldn't be invisible.",
    1, 3.8, 11, 0.6, size=20, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "github.com/revathi/herday",
    1, 5.2, 11, 0.5, size=16, color=LIGHT_PUR, align=PP_ALIGN.CENTER, italic=True)


prs.save("HerDay_Pitch.pptx")
print("HerDay_Pitch.pptx saved!")
